"""
agents/report.py — The Report Agent.

The FINAL agent in the pipeline. Runs after pitch.py.

What it does:
  1. Reads pitch_output (PitchOutput) from state
  2. Renders each of the 12 slide objects into a real .pptx file
  3. Uses python-pptx — pure Python, no Node.js, no external tools
  4. Writes the output path to state["final_report_path"]

Design approach:
  - Brand colors come from the pitch content (teal #0D9488 + dark #1E293B)
  - Consistent layout system: every slide has a colored left accent bar,
    a title zone, and a content zone
  - Slide types use distinct layouts:
      cover        → full dark background, centered hero text
      bullet slides → left accent bar + title + bullet rows
      matrix slide  → table with checkmarks
      financials    → big stat callout cards
      ask           → dark background, closing line hero

  python-pptx coordinate system: EMU (English Metric Units)
  1 inch = 914400 EMU. Use Inches() and Pt() helpers everywhere.
"""

import os
from datetime import datetime
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from state import AppState

load_dotenv()

# ── Output directory ──────────────────────────────────────────────────────────
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "outputs")


# ── Brand colors (no # prefix for python-pptx RGBColor) ─────────────────────
class C:
    TEAL        = RGBColor(0x0D, 0x94, 0x88)   # primary
    DARK        = RGBColor(0x1E, 0x29, 0x3B)   # dark slate
    AMBER       = RGBColor(0xF5, 0x9E, 0x0B)   # accent
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    OFF_WHITE   = RGBColor(0xF8, 0xFA, 0xFC)
    LIGHT_GRAY  = RGBColor(0xE2, 0xE8, 0xF0)
    MID_GRAY    = RGBColor(0x64, 0x74, 0x8B)
    NEAR_BLACK  = RGBColor(0x0F, 0x17, 0x2A)
    GREEN_CHECK = RGBColor(0x16, 0xA3, 0x4A)
    RED_CROSS   = RGBColor(0xDC, 0x26, 0x26)


# ── Slide dimensions: LAYOUT_WIDE 13.3" × 7.5" ───────────────────────────────
W = Inches(13.3)
H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# HELPER FUNCTIONS
# ─────────────────────────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    """Add a fully blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color: RGBColor):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    return shape


def _text(slide, text: str, x, y, w, h,
          size=18, bold=False, color=C.NEAR_BLACK,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _slide_label(slide, label: str):
    """Small ALL-CAPS category label top-right."""
    _text(slide, label.upper(),
          x=Inches(10.5), y=Inches(0.2), w=Inches(2.6), h=Inches(0.35),
          size=8, color=C.MID_GRAY, align=PP_ALIGN.RIGHT)


def _accent_bar(slide):
    """Left-side teal accent bar — the visual motif repeated on every content slide."""
    _rect(slide, Inches(0), Inches(0), Inches(0.12), H, C.TEAL)


def _title_zone(slide, title: str, subtitle: str = ""):
    """Standard title + optional subtitle in the top zone."""
    _text(slide, title,
          x=Inches(0.35), y=Inches(0.25), w=Inches(12.5), h=Inches(0.75),
          size=28, bold=True, color=C.NEAR_BLACK)
    if subtitle:
        _text(slide, subtitle,
              x=Inches(0.35), y=Inches(0.95), w=Inches(12.5), h=Inches(0.4),
              size=14, color=C.MID_GRAY, italic=True)


def _divider(slide, y):
    """Thin horizontal rule below title zone."""
    _rect(slide, Inches(0.35), y, Inches(12.6), Inches(0.02), C.LIGHT_GRAY)


def _bullet_row(slide, headline: str, supporting: str, y, icon: str = "●"):
    """
    One bullet row: colored dot + bold headline + supporting text below.
    Returns the y-position after this row.
    """
    row_h = Inches(0.9)

    # dot
    _rect(slide, Inches(0.35), y + Inches(0.12),
          Inches(0.06), Inches(0.06), C.TEAL)

    # headline
    _text(slide, headline,
          x=Inches(0.55), y=y, w=Inches(12.1), h=Inches(0.4),
          size=16, bold=True, color=C.NEAR_BLACK)

    # supporting
    _text(slide, supporting,
          x=Inches(0.55), y=y + Inches(0.38), w=Inches(12.1), h=Inches(0.45),
          size=12, color=C.MID_GRAY)

    return y + row_h


def _stat_card(slide, x, y, w, h, value: str, label: str,
               bg=C.TEAL, text_color=C.WHITE):
    """A big-number stat card."""
    _rect(slide, x, y, w, h, bg)
    # value
    _text(slide, value,
          x=x, y=y + Inches(0.15), w=w, h=Inches(0.65),
          size=28, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    # label
    _text(slide, label,
          x=x, y=y + Inches(0.75), w=w, h=Inches(0.4),
          size=10, color=text_color, align=PP_ALIGN.CENTER)


def _presenter_notes(slide, note: str):
    """Add presenter notes to the slide's notes pane."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = note


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE RENDERERS — one function per slide
# ─────────────────────────────────────────────────────────────────────────────

def _render_cover(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.DARK)

    # teal left bar (thicker on cover)
    _rect(slide, Inches(0), Inches(0), Inches(0.35), H, C.TEAL)

    # amber bottom strip
    _rect(slide, Inches(0), H - Inches(0.18), W, Inches(0.18), C.AMBER)

    # startup name — hero
    _text(slide, slide_data.startup_name,
          x=Inches(0.6), y=Inches(1.6), w=Inches(12), h=Inches(1.4),
          size=64, bold=True, color=C.WHITE, align=PP_ALIGN.LEFT)

    # tagline
    _text(slide, slide_data.tagline,
          x=Inches(0.6), y=Inches(3.0), w=Inches(10), h=Inches(0.6),
          size=22, bold=False, color=C.AMBER, align=PP_ALIGN.LEFT)

    # one-liner
    _text(slide, slide_data.one_liner,
          x=Inches(0.6), y=Inches(3.7), w=Inches(11), h=Inches(0.7),
          size=14, color=C.LIGHT_GRAY, italic=True, align=PP_ALIGN.LEFT)

    # year bottom right
    _text(slide, str(datetime.now().year),
          x=Inches(11.5), y=Inches(6.9), w=Inches(1.6), h=Inches(0.4),
          size=11, color=C.MID_GRAY, align=PP_ALIGN.RIGHT)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_problem(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.OFF_WHITE)
    _accent_bar(slide)
    _slide_label(slide, "The Problem")
    _title_zone(slide, slide_data.headline)
    _divider(slide, Inches(1.35))

    y = Inches(1.55)
    for bp in slide_data.pain_points:
        y = _bullet_row(slide, bp.headline, bp.supporting, y)

    # emotional hook — amber callout box
    _rect(slide, Inches(0.35), Inches(6.1), Inches(12.6), Inches(1.05), C.AMBER)
    _text(slide, slide_data.emotional_hook,
          x=Inches(0.55), y=Inches(6.2), w=Inches(12.2), h=Inches(0.8),
          size=13, bold=True, color=C.NEAR_BLACK, italic=True)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_solution(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.OFF_WHITE)
    _accent_bar(slide)
    _slide_label(slide, "The Solution")
    _title_zone(slide, slide_data.headline)
    _divider(slide, Inches(1.35))

    y = Inches(1.55)
    for bp in slide_data.solution_bullets:
        y = _bullet_row(slide, bp.headline, bp.supporting, y)

    # aha moment — teal callout
    _rect(slide, Inches(0.35), Inches(6.1), Inches(12.6), Inches(1.05), C.TEAL)
    _text(slide, slide_data.aha_moment,
          x=Inches(0.55), y=Inches(6.2), w=Inches(12.2), h=Inches(0.8),
          size=13, bold=True, color=C.WHITE, italic=True)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_product(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.OFF_WHITE)
    _accent_bar(slide)
    _slide_label(slide, "Product")
    _title_zone(slide, slide_data.headline)
    _divider(slide, Inches(1.35))

    # 4 feature cards in a 2×2 grid
    features = slide_data.core_features[:4]
    positions = [
        (Inches(0.35), Inches(1.55)),
        (Inches(6.85), Inches(1.55)),
        (Inches(0.35), Inches(3.6)),
        (Inches(6.85), Inches(3.6)),
    ]
    card_w, card_h = Inches(6.3), Inches(1.75)

    for i, (feat, (cx, cy)) in enumerate(zip(features, positions)):
        bg_color = C.TEAL if i % 2 == 0 else C.DARK
        _rect(slide, cx, cy, card_w, card_h, bg_color)
        _text(slide, feat,
              x=cx + Inches(0.15), y=cy + Inches(0.15),
              w=card_w - Inches(0.3), h=card_h - Inches(0.3),
              size=13, color=C.WHITE, bold=False, wrap=True)

    # demo flow strip at bottom
    _rect(slide, Inches(0.35), Inches(6.0), Inches(12.6), Inches(1.2), C.NEAR_BLACK)
    _text(slide, "Demo Flow  →  " + slide_data.demo_flow,
          x=Inches(0.55), y=Inches(6.1), w=Inches(12.2), h=Inches(0.95),
          size=11, color=C.LIGHT_GRAY, italic=True)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_market(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.OFF_WHITE)
    _accent_bar(slide)
    _slide_label(slide, "Market Size")
    _title_zone(slide, slide_data.headline)
    _divider(slide, Inches(1.35))

    # TAM / SAM / SOM stat cards
    card_w = Inches(3.9)
    card_h = Inches(2.2)
    card_y = Inches(1.55)
    gap    = Inches(0.25)

    cards = [
        ("TAM", slide_data.tam, C.TEAL),
        ("SAM", slide_data.sam, C.DARK),
        ("SOM", slide_data.som, C.AMBER),
    ]
    for i, (label, value, color) in enumerate(cards):
        cx = Inches(0.35) + i * (card_w + gap)
        _rect(slide, cx, card_y, card_w, card_h, color)
        _text(slide, label,
              x=cx, y=card_y + Inches(0.12), w=card_w, h=Inches(0.35),
              size=11, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        _text(slide, value,
              x=cx + Inches(0.1), y=card_y + Inches(0.45),
              w=card_w - Inches(0.2), h=card_h - Inches(0.55),
              size=11, color=C.WHITE, wrap=True)

    # tailwinds
    _text(slide, "Market Tailwinds",
          x=Inches(0.35), y=Inches(4.05), w=Inches(6), h=Inches(0.4),
          size=14, bold=True, color=C.NEAR_BLACK)

    ty = Inches(4.5)
    for t in slide_data.market_tailwinds:
        _rect(slide, Inches(0.35), ty + Inches(0.12),
              Inches(0.06), Inches(0.06), C.AMBER)
        _text(slide, t,
              x=Inches(0.55), y=ty, w=Inches(12.1), h=Inches(0.45),
              size=12, color=C.NEAR_BLACK)
        ty += Inches(0.52)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_business(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.OFF_WHITE)
    _accent_bar(slide)
    _slide_label(slide, "Business Model")
    _title_zone(slide, slide_data.headline)
    _divider(slide, Inches(1.35))

    # model description
    _text(slide, slide_data.model_description,
          x=Inches(0.35), y=Inches(1.5), w=Inches(12.6), h=Inches(0.7),
          size=13, color=C.MID_GRAY, italic=True)

    # pricing tiers
    tier_colors = [C.LIGHT_GRAY, C.TEAL, C.DARK]
    text_colors = [C.NEAR_BLACK, C.WHITE, C.WHITE]
    tier_w = Inches(4.0)
    tier_h = Inches(2.6)
    tier_y = Inches(2.35)

    for i, (tier, bg, tc) in enumerate(
        zip(slide_data.pricing_tiers, tier_colors, text_colors)
    ):
        cx = Inches(0.35) + i * (tier_w + Inches(0.2))
        _rect(slide, cx, tier_y, tier_w, tier_h, bg)
        _text(slide, tier,
              x=cx + Inches(0.15), y=tier_y + Inches(0.15),
              w=tier_w - Inches(0.3), h=tier_h - Inches(0.3),
              size=12, color=tc, wrap=True)

    # unit economics row
    _text(slide, "Unit Economics",
          x=Inches(0.35), y=Inches(5.15), w=Inches(6), h=Inches(0.35),
          size=13, bold=True, color=C.NEAR_BLACK)

    ue_w = Inches(4.0)
    for i, metric in enumerate(slide_data.unit_economics):
        cx = Inches(0.35) + i * (ue_w + Inches(0.2))
        _rect(slide, cx, Inches(5.55), ue_w, Inches(0.9), C.NEAR_BLACK)
        _text(slide, metric,
              x=cx + Inches(0.1), y=Inches(5.6),
              w=ue_w - Inches(0.2), h=Inches(0.75),
              size=12, bold=True, color=C.AMBER, align=PP_ALIGN.CENTER)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_traction(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.OFF_WHITE)
    _accent_bar(slide)
    _slide_label(slide, "Traction")
    _title_zone(slide, slide_data.headline)
    _divider(slide, Inches(1.35))

    y = Inches(1.55)
    for tp in slide_data.traction_points:
        y = _bullet_row(slide, tp.headline, tp.supporting, y)

    # user quote block
    _rect(slide, Inches(0.35), Inches(5.2), Inches(12.6), Inches(1.2), C.TEAL)
    _text(slide, f'"{slide_data.validation_quote}"',
          x=Inches(0.6), y=Inches(5.3), w=Inches(12.1), h=Inches(1.0),
          size=12, italic=True, color=C.WHITE)

    # next milestones compact row
    _text(slide, "Next Milestones:",
          x=Inches(0.35), y=Inches(6.55), w=Inches(3), h=Inches(0.4),
          size=11, bold=True, color=C.NEAR_BLACK)
    milestone_text = "  ·  ".join(slide_data.next_milestones)
    _text(slide, milestone_text,
          x=Inches(2.5), y=Inches(6.55), w=Inches(10.5), h=Inches(0.4),
          size=11, color=C.MID_GRAY)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_competition(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.OFF_WHITE)
    _accent_bar(slide)
    _slide_label(slide, "Competition")
    _title_zone(slide, slide_data.headline)
    _divider(slide, Inches(1.35))

    cols = ["WhatsApp", "GST", "Auto-Remind", "Freelancer", "₹ Affordable"]
    col_w  = Inches(1.9)
    row_h  = Inches(0.55)
    name_w = Inches(2.5)
    start_x = Inches(0.35)
    start_y = Inches(1.5)

    # header row
    _rect(slide, start_x, start_y, name_w, row_h, C.DARK)
    _text(slide, "Competitor", x=start_x, y=start_y,
          w=name_w, h=row_h, size=11, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
    for j, col in enumerate(cols):
        cx = start_x + name_w + j * col_w
        _rect(slide, cx, start_y, col_w, row_h, C.DARK)
        _text(slide, col, x=cx, y=start_y,
              w=col_w, h=row_h, size=10, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # data rows
    feature_keys = [
        "whatsapp_native", "gst_compliant", "auto_reminders",
        "freelancer_focused", "affordable_inr"
    ]
    for i, row in enumerate(slide_data.competitor_matrix):
        ry = start_y + (i + 1) * row_h
        is_us = row.is_us
        row_bg = C.TEAL if is_us else (C.OFF_WHITE if i % 2 == 0 else C.LIGHT_GRAY)
        name_color = C.WHITE if is_us else C.NEAR_BLACK

        _rect(slide, start_x, ry, name_w, row_h, row_bg)
        label = ("✦ " + row.competitor_name) if is_us else row.competitor_name
        _text(slide, label, x=start_x, y=ry,
              w=name_w, h=row_h, size=11, bold=is_us,
              color=name_color, align=PP_ALIGN.CENTER)

        values = [
            row.whatsapp_native, row.gst_compliant, row.auto_reminders,
            row.freelancer_focused, row.affordable_inr
        ]
        for j, val in enumerate(values):
            cx = start_x + name_w + j * col_w
            _rect(slide, cx, ry, col_w, row_h, row_bg)
            mark  = "✓" if val else "✗"
            mcolor = C.WHITE if is_us else (C.GREEN_CHECK if val else C.RED_CROSS)
            _text(slide, mark, x=cx, y=ry, w=col_w, h=row_h,
                  size=16, bold=True, color=mcolor, align=PP_ALIGN.CENTER)

    # our moat
    moat_y = start_y + (len(slide_data.competitor_matrix) + 1) * row_h + Inches(0.25)
    _rect(slide, start_x, moat_y, Inches(12.6), Inches(0.75), C.AMBER)
    _text(slide, "Our Moat:  " + slide_data.our_moat,
          x=start_x + Inches(0.15), y=moat_y + Inches(0.08),
          w=Inches(12.3), h=Inches(0.6),
          size=12, bold=True, color=C.NEAR_BLACK)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_gtm(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.OFF_WHITE)
    _accent_bar(slide)
    _slide_label(slide, "Go-To-Market")
    _title_zone(slide, slide_data.headline)
    _divider(slide, Inches(1.35))

    # 3-phase columns
    phases = [
        ("Phase 1 · 0→100", slide_data.phase_1, C.TEAL),
        ("Phase 2 · 100→1K", slide_data.phase_2, C.DARK),
        ("Phase 3 · 1K→10K", slide_data.phase_3, C.AMBER),
    ]
    col_w = Inches(4.1)
    col_h = Inches(3.0)
    col_y = Inches(1.55)

    for i, (label, text, color) in enumerate(phases):
        cx = Inches(0.35) + i * (col_w + Inches(0.2))
        _rect(slide, cx, col_y, col_w, Inches(0.4), color)
        _text(slide, label, x=cx, y=col_y, w=col_w, h=Inches(0.4),
              size=11, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        _rect(slide, cx, col_y + Inches(0.4), col_w, col_h - Inches(0.4), C.LIGHT_GRAY)
        _text(slide, text,
              x=cx + Inches(0.1), y=col_y + Inches(0.5),
              w=col_w - Inches(0.2), h=col_h - Inches(0.6),
              size=11, color=C.NEAR_BLACK, wrap=True)

    # channels + north star row at bottom
    _text(slide, "Primary Channels:",
          x=Inches(0.35), y=Inches(5.0), w=Inches(3), h=Inches(0.4),
          size=12, bold=True, color=C.NEAR_BLACK)
    channels_text = "  ·  ".join(slide_data.primary_channels)
    _text(slide, channels_text,
          x=Inches(2.8), y=Inches(5.0), w=Inches(10.2), h=Inches(0.4),
          size=12, color=C.MID_GRAY)

    _rect(slide, Inches(0.35), Inches(5.6), Inches(12.6), Inches(0.8), C.NEAR_BLACK)
    _text(slide, "⭐  " + slide_data.north_star,
          x=Inches(0.55), y=Inches(5.68), w=Inches(12.2), h=Inches(0.6),
          size=12, bold=True, color=C.AMBER)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_team(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.OFF_WHITE)
    _accent_bar(slide)
    _slide_label(slide, "Team")
    _title_zone(slide, slide_data.headline)
    _divider(slide, Inches(1.35))

    # why us block
    _rect(slide, Inches(0.35), Inches(1.55), Inches(12.6), Inches(2.0), C.DARK)
    _text(slide, slide_data.why_us,
          x=Inches(0.55), y=Inches(1.65), w=Inches(12.2), h=Inches(1.75),
          size=14, color=C.WHITE, italic=True, wrap=True)

    # key hires
    _text(slide, "Key Hires Needed",
          x=Inches(0.35), y=Inches(3.75), w=Inches(6), h=Inches(0.4),
          size=14, bold=True, color=C.NEAR_BLACK)

    hy = Inches(4.2)
    for hire in slide_data.key_hires_needed:
        _rect(slide, Inches(0.35), hy + Inches(0.1),
              Inches(0.06), Inches(0.06), C.TEAL)
        _text(slide, hire,
              x=Inches(0.55), y=hy, w=Inches(12.1), h=Inches(0.45),
              size=12, color=C.NEAR_BLACK)
        hy += Inches(0.52)

    # advisors
    _rect(slide, Inches(0.35), Inches(6.1), Inches(12.6), Inches(0.95), C.LIGHT_GRAY)
    _text(slide, "Advisors / Supporters:  " + slide_data.advisors_or_supporters,
          x=Inches(0.55), y=Inches(6.2), w=Inches(12.2), h=Inches(0.75),
          size=11, color=C.MID_GRAY, italic=True)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_financials(prs, slide_data, finance_output):
    slide = _blank_slide(prs)
    _bg(slide, C.OFF_WHITE)
    _accent_bar(slide)
    _slide_label(slide, "Financials")
    _title_zone(slide, slide_data.headline)
    _divider(slide, Inches(1.35))

    snap = slide_data.snapshot

    # 4 big stat cards top row
    stats = [
        (snap.arr_month_12,        "ARR · Month 12"),
        (snap.paid_users_month_12, "Paying Users · M12"),
        (snap.ltv_cac,             "LTV : CAC"),
        (snap.break_even_month,    "Break-Even"),
    ]
    card_w = Inches(3.0)
    card_h = Inches(1.5)
    gap    = Inches(0.18)
    card_y = Inches(1.55)
    card_colors = [C.TEAL, C.DARK, C.TEAL, C.DARK]

    for i, ((val, lbl), color) in enumerate(zip(stats, card_colors)):
        cx = Inches(0.35) + i * (card_w + gap)
        _stat_card(slide, cx, card_y, card_w, card_h, val, lbl, bg=color)

    # secondary stats row
    stats2 = [
        (snap.mrr_month_12, "MRR · Month 12"),
        (snap.gross_margin,  "Gross Margin"),
        (snap.runway,        "Runway"),
        (snap.raise_amount,  "Raising"),
    ]
    card_y2 = card_y + card_h + Inches(0.18)
    sec_colors = [C.AMBER, C.NEAR_BLACK, C.AMBER, C.NEAR_BLACK]
    for i, ((val, lbl), color) in enumerate(zip(stats2, sec_colors)):
        cx = Inches(0.35) + i * (card_w + gap)
        _stat_card(slide, cx, card_y2, card_w, Inches(1.2), val, lbl, bg=color)

    # projection narrative
    _text(slide, slide_data.projection_narrative,
          x=Inches(0.35), y=Inches(4.8), w=Inches(12.6), h=Inches(0.75),
          size=12, color=C.MID_GRAY, italic=True)

    # assumptions strip
    _rect(slide, Inches(0.35), Inches(5.7), Inches(12.6), Inches(0.6), C.LIGHT_GRAY)
    assumptions_text = "  ·  ".join(slide_data.key_assumptions)
    _text(slide, "Assumptions:  " + assumptions_text,
          x=Inches(0.55), y=Inches(5.78), w=Inches(12.2), h=Inches(0.45),
          size=10, color=C.MID_GRAY)

    _presenter_notes(slide, slide_data.presenter_note)


def _render_ask(prs, slide_data):
    slide = _blank_slide(prs)
    _bg(slide, C.DARK)

    # teal left bar
    _rect(slide, Inches(0), Inches(0), Inches(0.35), H, C.TEAL)

    # amber bottom strip
    _rect(slide, Inches(0), H - Inches(0.18), W, Inches(0.18), C.AMBER)

    _slide_label(slide, "The Ask")

    # headline
    _text(slide, slide_data.headline,
          x=Inches(0.55), y=Inches(0.3), w=Inches(12.5), h=Inches(0.8),
          size=24, bold=True, color=C.WHITE)

    # raise amount big
    _text(slide, slide_data.raise_amount,
          x=Inches(0.55), y=Inches(1.1), w=Inches(12.5), h=Inches(0.75),
          size=36, bold=True, color=C.AMBER)

    # use of funds — 4 cards
    uof_w = Inches(3.0)
    uof_h = Inches(1.5)
    uof_y = Inches(2.0)
    uof_colors = [C.TEAL, C.NEAR_BLACK, C.TEAL, C.NEAR_BLACK]

    for i, (item, color) in enumerate(zip(slide_data.use_of_funds, uof_colors)):
        cx = Inches(0.55) + i * (uof_w + Inches(0.22))
        _rect(slide, cx, uof_y, uof_w, uof_h, color)
        _text(slide, item,
              x=cx + Inches(0.1), y=uof_y + Inches(0.1),
              w=uof_w - Inches(0.2), h=uof_h - Inches(0.2),
              size=11, color=C.WHITE, wrap=True)

    # milestones unlocked
    _text(slide, "Milestones This Unlocks",
          x=Inches(0.55), y=Inches(3.7), w=Inches(6), h=Inches(0.4),
          size=13, bold=True, color=C.LIGHT_GRAY)

    my = Inches(4.15)
    for m in slide_data.milestones_unlocked:
        _rect(slide, Inches(0.55), my + Inches(0.1),
              Inches(0.06), Inches(0.06), C.AMBER)
        _text(slide, m, x=Inches(0.75), y=my,
              w=Inches(12.1), h=Inches(0.45),
              size=12, color=C.LIGHT_GRAY)
        my += Inches(0.5)

    # closing line — hero
    _rect(slide, Inches(0.35), Inches(5.8), Inches(12.6), Inches(1.4), C.TEAL)
    _text(slide, slide_data.closing_line,
          x=Inches(0.55), y=Inches(5.9), w=Inches(12.2), h=Inches(1.2),
          size=16, bold=True, color=C.WHITE, italic=True,
          align=PP_ALIGN.CENTER)

    _presenter_notes(slide, slide_data.presenter_note)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN AGENT FUNCTION
# ─────────────────────────────────────────────────────────────────────────────

def run_report_agent(state: AppState) -> AppState:
    """
    Main entry point. Reads pitch_output from state.
    Renders a full 12-slide PPTX and writes the path to state.
    """

    print("\n[Report Agent] Rendering pitch deck...")

    if not state.get("pitch_output"):
        errors = state.get("errors") or []
        errors.append("report_agent: pitch_output is missing — run pitch agent first")
        return {**state, "errors": errors}

    pitch   = state["pitch_output"]
    finance = state.get("finance_output")

    # ── Set up presentation ───────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # ── Render all 12 slides in order ─────────────────────────────────────
    _render_cover      (prs, pitch.slide_01_cover)
    _render_problem    (prs, pitch.slide_02_problem)
    _render_solution   (prs, pitch.slide_03_solution)
    _render_product    (prs, pitch.slide_04_product)
    _render_market     (prs, pitch.slide_05_market)
    _render_business   (prs, pitch.slide_06_business)
    _render_traction   (prs, pitch.slide_07_traction)
    _render_competition(prs, pitch.slide_08_competition)
    _render_gtm        (prs, pitch.slide_09_gtm)
    _render_team       (prs, pitch.slide_10_team)
    _render_financials (prs, pitch.slide_11_financials, finance)
    _render_ask        (prs, pitch.slide_12_ask)

    # ── Save file ─────────────────────────────────────────────────────────
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    safe_name = pitch.slide_01_cover.startup_name.replace(" ", "_").lower()
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    filename  = f"{safe_name}_pitch_deck_{timestamp}.pptx"
    filepath  = os.path.join(OUTPUT_DIR, filename)

    prs.save(filepath)
    print(f"[Report Agent] ✓ Deck saved → {filepath}")
    print(f"[Report Agent]   Slides: {len(prs.slides)}")
    print(f"[Report Agent]   File:   {filename}")

    # ── Write back to state ───────────────────────────────────────────────
    completed = state.get("completed_agents") or []
    completed.append("report")

    return {
        **state,
        "final_report_path": filepath,
        "completed_agents":  completed
    }